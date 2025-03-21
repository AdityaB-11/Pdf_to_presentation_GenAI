import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Define themes with specific files
THEMES = {
    'theme1': {
        'name': 'Theme 1',
        'file': 'Presentation1.pptx'
    },
    'theme2': {
        'name': 'Theme 2',
        'file': 'Presentation2.pptx'
    },
    'theme3': {
        'name': 'Theme 3',
        'file': 'Presentation3.pptx'
    },
    'theme4': {
        'name': 'Theme 4',
        'file': 'Presentation4.pptx'
    }
}

def list_available_themes():
    """Return a list of available themes."""
    return [{'key': key, 'name': theme['name']} for key, theme in THEMES.items()]

def parse_vba_file(vba_file_path):
    slides_data = []
    current_slide = None
    content_buffer = ""

    with open(vba_file_path, 'r', encoding='utf-8') as file:
        for line_number, line in enumerate(file, 1):
            line = line.strip()
            try:
                if "Set sld = ppt.Slides.Add" in line:
                    if current_slide:
                        current_slide['content'] = content_buffer.strip()
                        slides_data.append(current_slide)
                    current_slide = {'title': '', 'content': ''}
                    content_buffer = ""
                elif ".Shapes.Title.TextFrame.TextRange.Text =" in line:
                    match = re.search(r'"([^"]*)"', line)
                    if match:
                        current_slide['title'] = match.group(1)
                elif ".Text =" in line:
                    match = re.search(r'"([^"]*)"', line)
                    if match:
                        content_buffer += match.group(1)
                elif '& _' in line:
                    match = re.search(r'"([^"]*)"', line)
                    if match:
                        content_buffer += match.group(1) + "\n"
            except ValueError as e:
                print(f"Error processing line {line_number}: {line}")
                print(f"Error details: {str(e)}")

    if current_slide:
        current_slide['content'] = content_buffer.strip()
        slides_data.append(current_slide)

    return slides_data

def detect_format_type(slides_data):
    """Detect the format type based on the content structure"""
    # Look for indicators in the content
    format_indicators = {
        'qa_format': 0,
        'table_format': 0,
        'diagram_format': 0,
        'bullet_only': 0
    }
    
    # Check slide titles and content for format indicators
    for slide in slides_data:
        # Check for Q&A format
        if slide['title'].endswith('?'):
            format_indicators['qa_format'] += 1
            
        content = slide['content']
        # Check for table descriptions
        if 'table' in content.lower() and ('row' in content.lower() or 'column' in content.lower()):
            format_indicators['table_format'] += 1
            
        # Check for diagram descriptions
        if any(term in content.lower() for term in ['diagram', 'chart', 'graph', 'visualize']):
            format_indicators['diagram_format'] += 1
            
        # Check if content is short bullet points only
        if isinstance(content, list) and all(len(point) < 120 for point in content):
            format_indicators['bullet_only'] += 1
            
    # Determine the dominant format type
    if format_indicators['qa_format'] > len(slides_data) / 3:
        return 'qa'
    elif format_indicators['table_format'] > len(slides_data) / 3:
        return 'tables'
    elif format_indicators['diagram_format'] > len(slides_data) / 3:
        return 'diagrams'
    elif format_indicators['bullet_only'] > len(slides_data) / 2:
        return 'bullets_only'
    else:
        return 'standard'

def apply_theme_to_slide(slide, theme):
    """Apply theme colors and styles to a slide."""
    # Apply theme to title
    if slide.shapes.title:
        title_frame = slide.shapes.title.text_frame
        title_frame.paragraphs[0].font.color.rgb = theme['title_color']
        title_frame.paragraphs[0].font.size = theme['font_size']['title']
    
    # Apply theme to body
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if shape != slide.shapes.title:  # Skip title
                    paragraph.font.color.rgb = theme['body_color']
                    paragraph.font.size = theme['font_size']['body']

def get_image_titles():
    """Read image titles from the file."""
    titles_file = os.path.join('extract', 'image_titles.txt')
    titles = {}
    if os.path.exists(titles_file):
        with open(titles_file, 'r', encoding='utf-8') as f:
            for line in f:
                key, title = line.strip().split('|', 1)
                titles[key] = title
    return titles

def create_image_slide(prs, images_dir, page_number):
    """Create a new slide specifically for images."""
    # Check if images directory exists
    if not os.path.exists(images_dir):
        print(f"Images directory not found: {images_dir}")
        return None
        
    # Find all images for this page - check both PDF extraction format and topic generator format
    slide_images = [f for f in os.listdir(images_dir) 
                   if (f'page_{page_number}_img_' in f.lower() or f'topic_{page_number}_img_' in f.lower()) and 
                   any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
    
    if not slide_images:
        return None

    # Get image titles from the titles file
    image_titles = get_image_titles()
    
    # Sort images by their numbers - extract the number from different formats
    def get_img_number(filename):
        if 'img_' in filename:
            match = re.search(r'img_(\d+)', filename)
            if match:
                return int(match.group(1))
        return 0
    
    slide_images.sort(key=get_img_number)
    
    # Check for extracted text for this page to get better title
    page_text = ""
    extracted_text_path = os.path.join('extract', f'pdf_page_{page_number}.txt')
    if not os.path.exists(extracted_text_path):
        # Try alternative naming patterns
        pdf_files = [f for f in os.listdir('extract') if f.endswith(f'_page_{page_number}.txt')]
        if pdf_files:
            extracted_text_path = os.path.join('extract', pdf_files[0])
    
    if os.path.exists(extracted_text_path):
        with open(extracted_text_path, 'r', encoding='utf-8') as f:
            page_text = f.read()
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Try to extract a meaningful title from the page text or image titles
    title_text = ""
    
    # First, check image titles for good titles
    for img_file in slide_images:
        # Handle both PDF extraction and topic generator formats
        if 'page_' in img_file:
            img_index = re.search(r'img_(\d+)', img_file).group(1)
            img_key = f"page_{page_number}_img_{img_index}"
        elif 'topic_' in img_file:
            img_index = re.search(r'img_(\d+)', img_file).group(1)
            img_key = f"page_{page_number}_img_{img_index}"
        else:
            continue
            
        if img_key in image_titles and image_titles[img_key] and len(image_titles[img_key]) > 5:
            # This seems like a real title, not just "Figure X"
            if not image_titles[img_key].lower().startswith("figure"):
                title_text = image_titles[img_key]
                break
    
    # If no good title from image titles, try to extract from page text
    if not title_text and page_text:
        # Look for section titles or headers in the text
        title_patterns = [
            r'(?:^|\n)((?:[A-Z][a-z]*\s*){1,6}(?:Figures|Figure|Images|Charts|Diagrams|Illustrations))(?:\n|:)',
            r'(?:^|\n)(\d+\.\d+\s+(?:[A-Z][a-z]*\s*){1,6})(?:\n|:)',  # Section numbers like 2.1 Title
            r'(?:^|\n)((?:[A-Z][a-z]*\s*){2,6})(?:\n)',  # Capitalized phrases
        ]
        
        for pattern in title_patterns:
            matches = re.finditer(pattern, page_text)
            for match in matches:
                candidate = match.group(1).strip()
                if 4 < len(candidate) < 60:  # Reasonable title length
                    title_text = candidate
                    break
            if title_text:
                break
    
    # For topic generator, if we still don't have a good title, use a topic-related title
    if not title_text and any('topic_' in img for img in slide_images):
        title_text = "Visual Content"
        # Try to extract topic from image titles
        for img_key, img_title in image_titles.items():
            if not img_title.startswith("Figure"):
                title_text = f"Visual Content: {img_title}"
                break
                
    # If still no good title, try specific image-related text
    if not title_text and page_text:
        # Look for image references in text
        img_ref_patterns = [
            r'(?:Figure|Fig\.|Image|Diagram|Chart)\s*\d+[:\.\s]+([^\n\.]{10,60})',
            r'(?:see|in|the)\s+(?:figure|image|diagram|chart|illustration)\s+([^\n\.]{10,60})'
        ]
        
        for pattern in img_ref_patterns:
            matches = re.finditer(pattern, page_text, re.IGNORECASE)
            for match in matches:
                candidate = match.group(1).strip()
                if candidate:
                    title_text = candidate
                    break
            if title_text:
                break
                
    # Final fallback if we still don't have a good title
    if not title_text:
        # Try to get the first sentence or phrase from the page
        if page_text:
            lines = [line.strip() for line in page_text.split('\n') if line.strip()]
            for line in lines:
                if 10 < len(line) < 100 and not line.startswith('•'):
                    title_text = line
                    break
        
        # If still nothing, use the default
        if not title_text:
            title_text = f"Visual Content for Section {page_number}"
    
    # Add the title to the slide
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Calculate layout for images
    if len(slide_images) == 1:
        left_margin = Inches(1)
        top_margin = Inches(1.5)
        img_width = Inches(8)
        img_height = Inches(5.5)
    else:
        left_margin = Inches(0.5)
        top_margin = Inches(1.5)
        img_width = Inches(4.5)
        img_height = Inches(3.5)
        
    # Add images with captions
    for idx, image_file in enumerate(slide_images):
        image_path = os.path.join(images_dir, image_file)
        try:
            row = idx // 2
            col = idx % 2
            left = left_margin + (col * (img_width + Inches(0.5)))
            top = top_margin + (row * (img_height + Inches(0.7)))  # Extra space for caption
            
            # Add the image
            pic = slide.shapes.add_picture(
                image_path,
                left=left,
                top=top,
                width=img_width,
                height=img_height
            )
            
            # Add a caption below the image
            # Handle both formats
            if 'page_' in image_file:
                img_index = re.search(r'img_(\d+)', image_file).group(1)
                img_key = f"page_{page_number}_img_{img_index}"
            elif 'topic_' in image_file:
                img_index = re.search(r'img_(\d+)', image_file).group(1)
                img_key = f"page_{page_number}_img_{img_index}"
            else:
                img_key = None
                
            caption_text = f"Figure {idx+1}"  # Default caption
            if img_key and img_key in image_titles:
                caption_text = image_titles[img_key]
                
            caption = slide.shapes.add_textbox(
                left=left, 
                top=top + img_height + Inches(0.1),
                width=img_width,
                height=Inches(0.5)
            )
            caption_frame = caption.text_frame
            caption_frame.word_wrap = True
            caption_frame.text = caption_text
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            caption_frame.paragraphs[0].font.size = Pt(10)
            caption_frame.paragraphs[0].font.italic = True
            
            print(f"Added image: {image_file} with caption: {caption_text}")
        except Exception as e:
            print(f"Error adding image {image_file}: {str(e)}")
    
    return slide

def create_bullet_slide(prs, slide_title, bullet_points, format_type='standard'):
    """Create a slide with properly formatted bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    title = slide.shapes.title
    title.text = slide_title
    
    # Get content placeholder
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Format bullet points based on format type
    for i, point in enumerate(bullet_points):
        point_text = point.strip()
        if point_text.startswith('- '):
            point_text = point_text[2:]
            
        p = tf.add_paragraph()
        p.text = point_text
        p.level = 0
        
        # Apply special formatting based on format type
        if format_type == 'qa':
            # Format first bullet differently if it's a question
            if i == 0 and point_text.endswith('?'):
                p.font.bold = True
                p.font.size = Pt(18)
                
        elif format_type == 'tables' and 'table:' in point_text.lower():
            # Format table descriptions with special styling
            p.font.italic = True
        
        elif format_type == 'diagrams' and any(term in point_text.lower() for term in ['diagram:', 'chart:', 'graph:']):
            # Format diagram descriptions
            p.font.italic = True
            p.font.color.rgb = RGBColor(0, 112, 192)  # Blue color
    
    return slide

def create_qa_slide(prs, question, answers):
    """Create a slide in Question & Answer format"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set question as the title
    title = slide.shapes.title
    title.text = question
    
    # Add answers as bullet points
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    for answer in answers:
        answer_text = answer.strip()
        if answer_text.startswith('- '):
            answer_text = answer_text[2:]
            
        p = tf.add_paragraph()
        p.text = answer_text
        p.level = 0
    
    return slide

def create_diagram_slide(prs, slide_title, diagram_description, page_number=None):
    """Create a slide with a placeholder for a diagram"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = slide_title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Extract diagram type and description
    diagram_type = "Diagram"
    diagram_content = diagram_description
    
    # Try to identify diagram type
    diagram_types = ["flowchart", "pie chart", "bar chart", "line graph", "org chart", 
                    "venn diagram", "timeline", "mind map", "scatter plot", "diagram"]
    
    for d_type in diagram_types:
        if d_type in diagram_description.lower():
            diagram_type = d_type.title()
            break
    
    # Look for images that might match this diagram description
    images_dir = os.path.join('extract', 'images')
    diagram_images = []
    
    # Only look for specific page images if we have a page number
    if page_number and os.path.exists(images_dir):
        # Look specifically for images from this page, supporting both PDF extraction and topic generator
        page_images = [f for f in os.listdir(images_dir) 
                      if (f'page_{page_number}_img_' in f.lower() or f'topic_{page_number}_img_' in f.lower()) and 
                      any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
        
        # Get image titles to help find the most relevant one
        image_titles = get_image_titles()
        
        # First try: Find images with titles matching the diagram type
        for img_file in page_images:
            if 'img_' in img_file:
                match = re.search(r'img_(\d+)', img_file)
                if match:
                    img_index = match.group(1)
                    
                    # Handle both formats
                    if 'page_' in img_file:
                        img_key = f"page_{page_number}_img_{img_index}"
                    elif 'topic_' in img_file:
                        img_key = f"page_{page_number}_img_{img_index}"
                    else:
                        continue
                        
                    if img_key in image_titles:
                        title = image_titles[img_key].lower()
                        # Check if the image title contains diagram type keywords
                        if any(d_type.lower() in title for d_type in diagram_types):
                            diagram_images.append(img_file)
                        # Also check if the diagram description contains words from the image title
                        elif any(word in diagram_description.lower() for word in title.split() if len(word) > 3):
                            diagram_images.append(img_file)
        
        # If we didn't find any matching images, just use the first image from this page
        if not diagram_images and page_images:
            diagram_images = [page_images[0]]
    
    # If no page-specific images found, look for any images
    if not diagram_images and os.path.exists(images_dir):
        all_images = [f for f in os.listdir(images_dir) 
                    if any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
        if all_images:
            # Just use the first available image
            diagram_images = [all_images[0]]
    
    # If we have images, use the first one
    if diagram_images:
        try:
            # Add the image in the center
            image_path = os.path.join(images_dir, diagram_images[0])
            slide.shapes.add_picture(
                image_path,
                left=Inches(1),
                top=Inches(2),
                width=Inches(8),
                height=Inches(4)
            )
            
            # Get caption for the image
            img_caption = ""
            if page_number:
                img_file = diagram_images[0]
                if 'img_' in img_file:
                    match = re.search(r'img_(\d+)', img_file)
                    if match:
                        img_index = match.group(1)
                        
                        # Handle both formats
                        if 'page_' in img_file:
                            img_key = f"page_{page_number}_img_{img_index}"
                        elif 'topic_' in img_file:
                            img_key = f"page_{page_number}_img_{img_index}"
                        else:
                            img_key = None
                            
                        image_titles = get_image_titles()
                        if img_key and img_key in image_titles:
                            img_caption = image_titles[img_key]
            
            # Add description below with caption if available
            desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
            desc_frame = desc_box.text_frame
            desc_frame.word_wrap = True
            
            p = desc_frame.add_paragraph()
            # Use image caption if available, otherwise use the diagram description
            display_text = img_caption if img_caption else diagram_content
            p.text = display_text
            p.font.size = Pt(12)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
            
            print(f"Added diagram slide with image: {diagram_images[0]}")
            return slide
        except Exception as e:
            print(f"Error adding diagram image: {str(e)}")
    
    # If no images or error, add placeholder
    diagram_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    diagram_frame = diagram_box.text_frame
    diagram_frame.word_wrap = True
    
    p = diagram_frame.add_paragraph()
    p.text = f"[{diagram_type} Placeholder]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # Add description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    
    p = desc_frame.add_paragraph()
    p.text = diagram_content
    p.font.size = Pt(12)
    p.font.italic = True
    
    return slide

def create_powerpoint(slides_data, output_file, theme_key='theme1'):
    """Create PowerPoint with specified theme."""
    print(f"Creating PowerPoint with theme: {theme_key}")
    
    # Get theme configuration
    theme = THEMES.get(theme_key, THEMES['theme1'])
    theme_path = os.path.join('themes', theme['file'])
    
    # Create presentation with selected theme
    prs = Presentation(theme_path)
    
    # Detect format type from the content
    format_type = detect_format_type(slides_data)
    print(f"Detected presentation format: {format_type}")
    
    # Get the images directory
    images_dir = os.path.join('extract', 'images')
    has_images = os.path.exists(images_dir)
    if has_images:
        print(f"Looking for images in: {images_dir}")
    else:
        print(f"No images directory found at: {images_dir}. Skipping image slides.")
    
    # Check if we're in topic generator mode (no PDF source)
    topic_generator_mode = False
    if has_images:
        # Check if the images are from topic generator (named topic_X_img_Y.jpg) rather than PDF extraction
        topic_images = [f for f in os.listdir(images_dir) 
                       if 'topic_' in f.lower() and 
                       any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
        if topic_images:
            topic_generator_mode = True
            print("Detected topic generator mode with web images")
    
    # Preprocess image information
    image_titles = get_image_titles() if has_images else {}
    
    # Map of slide numbers to pages (for topic generator mode)
    slide_to_page_map = {}
    
    # First pass: map slide numbers to pages for topic generator mode
    if topic_generator_mode:
        for slide_num, slide_data in enumerate(slides_data, 1):
            # Skip title and agenda slides (usually first two)
            if slide_num <= 2:
                continue
                
            # For content slides, map to page numbers
            page_num = slide_num - 2  # Page 1 = slide 3 (after title and agenda)
            slide_to_page_map[slide_num] = page_num
    
    # Create content slides
    created_pages = set()  # Track which pages have been processed
    
    for slide_num, slide_data in enumerate(slides_data, 1):
        # Calculate corresponding page number
        if topic_generator_mode:
            page_number = slide_to_page_map.get(slide_num)
        else:
            page_number = slide_num - 2 if slide_num > 2 else None
        
        slide_title = slide_data['title']
        
        # Handle the content based on format type
        if isinstance(slide_data['content'], list):
            bullet_points = slide_data['content']
        else:
            # Split content into bullet points if it's a string
            bullet_points = [line.strip() for line in slide_data['content'].split('\n') if line.strip()]
        
        # Skip empty slides
        if not bullet_points and slide_num > 2:  # Allow empty title or index slide
            continue
            
        # Check if this is specifically a diagram slide
        is_diagram_slide = False
        if format_type == 'diagrams' and any(any(term in bp.lower() for term in ['diagram:', 'chart:', 'graph:']) for bp in bullet_points):
            is_diagram_slide = True
            diagram_desc = '\n'.join(bullet_points)
            
            # If it's a diagram slide and we have images, create a special slide
            if has_images and page_number and page_number > 0:
                # Find any images for this page - support both PDF and topic generator formats
                if topic_generator_mode:
                    page_images = [f for f in os.listdir(images_dir) 
                                  if (f'topic_{page_number}_img_' in f.lower()) and 
                                  any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
                else:
                    page_images = [f for f in os.listdir(images_dir) 
                                  if (f'page_{page_number}_img_' in f.lower()) and 
                                  any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
                
                if page_images:
                    # If there are images, use them on this diagram slide
                    create_diagram_slide(prs, slide_title, diagram_desc, page_number)
                    print(f"Created diagram slide with embedded image: {slide_title}")
                    created_pages.add(page_number)
                    continue
            
        # Create the appropriate type of slide based on format
        if format_type == 'qa' and slide_title.endswith('?'):
            create_qa_slide(prs, slide_title, bullet_points)
            print(f"Created Q&A slide: {slide_title}")
        elif is_diagram_slide:
            create_diagram_slide(prs, slide_title, diagram_desc, page_number)
            print(f"Created diagram slide: {slide_title}")
        else:
            # Default to bullet point slide
            create_bullet_slide(prs, slide_title, bullet_points, format_type)
            print(f"Created standard slide: {slide_title}")
        
        # For topic generator mode, try to insert relevant images after key content slides
        if topic_generator_mode and has_images and page_number and page_number > 0:
            # Only add image slide if we haven't already used this page number
            if page_number not in created_pages:
                # Every other slide, check if we have images for this "page"
                if page_number % 2 == 0:  # Even page numbers get images
                    page_images = [f for f in os.listdir(images_dir) 
                                  if f'topic_{page_number}_img_' in f.lower() and 
                                  any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
                    
                    if page_images:
                        image_slide = create_image_slide(prs, images_dir, page_number)
                        if image_slide:
                            print(f"Added web image slide after content slide {slide_num}")
                            created_pages.add(page_number)
        
        # Add image slide for PDF extraction (non-diagram case)
        elif not topic_generator_mode and has_images and page_number and page_number > 0 and not is_diagram_slide:
            # Only add image slide if we haven't already used this page number
            if page_number not in created_pages:
                # Check if there are images for this page number first
                page_images = [f for f in os.listdir(images_dir) 
                              if f'page_{page_number}_img_' in f.lower() and 
                              any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif'])]
                    
                if page_images:
                    image_slide = create_image_slide(prs, images_dir, page_number)
                    if image_slide:
                        print(f"Added image slide for page {page_number}")
                        created_pages.add(page_number)
    
    # For topic generator mode, add any unused images at the end
    if topic_generator_mode and has_images:
        # Find all unused topic images
        all_topic_images = {}
        for img_file in os.listdir(images_dir):
            if 'topic_' in img_file.lower() and any(img_file.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif']):
                match = re.search(r'topic_(\d+)_img_', img_file.lower())
                if match:
                    page_num = int(match.group(1))
                    if page_num not in created_pages:
                        if page_num not in all_topic_images:
                            all_topic_images[page_num] = []
                        all_topic_images[page_num].append(img_file)
        
        # Add slides for any remaining images
        for page_num in sorted(all_topic_images.keys()):
            if all_topic_images[page_num]:
                image_slide = create_image_slide(prs, images_dir, page_num)
                if image_slide:
                    print(f"Added additional web image slide for content {page_num}")

    prs.save(output_file)
    print(f"Presentation saved with theme {theme_key} to: {output_file}")
    return output_file

def main():
    vba_file_path = 'create_presentation.vba'
    output_file = os.path.join('output', 'generated_presentation.pptx')
    
    # List available themes
    print("\nAvailable themes:")
    for theme in list_available_themes():
        print(f"- {theme['name']} (key: {theme['key']})")
    
    slides_data = parse_vba_file(vba_file_path)
    ppt_path = create_powerpoint(slides_data, output_file)
    return ppt_path

if __name__ == "__main__":
    main()
